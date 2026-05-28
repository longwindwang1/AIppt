"""Deck (Pydantic) → .pptx 文件。

M10 重构：彻底改视觉
- 大字号（bullets 28 / title bar 32 / cover 56）
- Cover 左侧整条色带 + 右侧大字 + chip 标签
- Section 大编号 + 居中标题
- Example 分 3 个色块区（题目 / 解答 / 答案）
- Practice 题目卡片化 + 留白答题框
- Interactive 大字居中问题
- Summary 编号要点 + 视觉强调
- 每页 footer：页码 + 课题（从 deck 传入）
- diagram 占满宽度 + 标题
"""
from __future__ import annotations

import io
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.config import settings
from app.models.slide import AnimationHint, Deck, Slide, SlideType
from app.services import diagrams


# --- 主题系统 ---------------------------------------------------------------

def _rgb_to_hex(c: RGBColor | None) -> str:
    if c is None:
        return "#ffffff"
    return "#" + bytes(c).hex()


def _lighten(c: RGBColor, ratio: float) -> RGBColor:
    """把颜色按 ratio (0~1) 向白色靠近。ratio=0.9 = 几乎白。"""
    r, g, b = bytes(c)[0], bytes(c)[1], bytes(c)[2]
    nr = int(r + (255 - r) * ratio)
    ng = int(g + (255 - g) * ratio)
    nb = int(b + (255 - b) * ratio)
    return RGBColor(nr, ng, nb)


def _darken(c: RGBColor, ratio: float) -> RGBColor:
    r, g, b = bytes(c)[0], bytes(c)[1], bytes(c)[2]
    return RGBColor(int(r * (1 - ratio)), int(g * (1 - ratio)), int(b * (1 - ratio)))


@dataclass(frozen=True)
class Theme:
    key: str
    label: str
    cn_font: str
    title_color: RGBColor
    accent_color: RGBColor
    text_color: RGBColor
    muted_color: RGBColor
    bg_color: RGBColor | None = None
    title_bar_text_color: RGBColor = field(default_factory=lambda: RGBColor(0xFF, 0xFF, 0xFF))
    base_size_bullet: int = 28              # M10：22 → 28
    base_size_title_bar: int = 32           # M10：28 → 32
    base_size_cover_title: int = 56         # M10：44 → 56
    rounded_corners: bool = False

    def best_for(self) -> str:
        return {
            "formal_blue": "公开课、家长会，正式场合",
            "kids_warm": "低年级日常课，活泼亲切",
            "blackboard": "复习课，黑板感强",
        }.get(self.key, "")

    @property
    def title_hex(self) -> str: return _rgb_to_hex(self.title_color)
    @property
    def accent_hex(self) -> str: return _rgb_to_hex(self.accent_color)
    @property
    def text_hex(self) -> str: return _rgb_to_hex(self.text_color)
    @property
    def muted_hex(self) -> str: return _rgb_to_hex(self.muted_color)
    @property
    def bg_hex(self) -> str: return _rgb_to_hex(self.bg_color) if self.bg_color else "#ffffff"
    @property
    def title_fg_hex(self) -> str: return _rgb_to_hex(self.title_bar_text_color)

    # M10：派生色 — soft 背景（区块底色）和 deep 强调色
    @property
    def soft_title(self) -> RGBColor:
        if self.bg_color is not None:
            return _darken(self.bg_color, -0.15) if False else RGBColor(0x2A, 0x3D, 0x2E)
        return _lighten(self.title_color, 0.92)

    @property
    def soft_accent(self) -> RGBColor:
        if self.bg_color is not None:
            return RGBColor(0x3D, 0x2A, 0x1E)
        return _lighten(self.accent_color, 0.90)


_FORMAL_BLUE = Theme(
    key="formal_blue",
    label="正式蓝",
    cn_font="Microsoft YaHei",
    title_color=RGBColor(0x1F, 0x4E, 0x79),
    accent_color=RGBColor(0xE8, 0x74, 0x22),
    text_color=RGBColor(0x2C, 0x3E, 0x50),
    muted_color=RGBColor(0x7F, 0x8C, 0x8D),
)

_KIDS_WARM = Theme(
    key="kids_warm",
    label="童趣暖",
    cn_font="Microsoft YaHei",
    title_color=RGBColor(0xE6, 0x7E, 0x22),
    accent_color=RGBColor(0x27, 0xAE, 0x60),
    text_color=RGBColor(0x34, 0x49, 0x5E),
    muted_color=RGBColor(0x95, 0xA5, 0xA6),
    base_size_bullet=30,
    base_size_title_bar=34,
    base_size_cover_title=60,
    rounded_corners=True,
)

_BLACKBOARD = Theme(
    key="blackboard",
    label="黑板风",
    cn_font="Microsoft YaHei",
    title_color=RGBColor(0xF1, 0xC4, 0x0F),
    accent_color=RGBColor(0xE7, 0x4C, 0x3C),
    text_color=RGBColor(0xEC, 0xF0, 0xF1),
    muted_color=RGBColor(0xBD, 0xC3, 0xC7),
    bg_color=RGBColor(0x1B, 0x2D, 0x1F),
    title_bar_text_color=RGBColor(0xF1, 0xC4, 0x0F),
)


THEMES: dict[str, Theme] = {t.key: t for t in [_FORMAL_BLUE, _KIDS_WARM, _BLACKBOARD]}


def get_theme(key: str | None) -> Theme:
    if not key or key not in THEMES:
        return _FORMAL_BLUE
    return THEMES[key]


# --- 公开 API ---------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def render(deck: Deck, out_path: Path | None = None,
           theme_key: str | None = None) -> Path:
    theme = get_theme(theme_key)
    prs = _new_presentation()

    expanded = _expand_animations(deck.slides)
    total = len(expanded)
    for i, slide in enumerate(expanded):
        _add_slide(prs, slide, theme, page_num=i + 1, total=total, lesson=deck.lesson_name)

    if out_path is None:
        settings.runs_dir.mkdir(parents=True, exist_ok=True)
        out_path = settings.runs_dir / deck.filename()

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def render_to_bytes(deck: Deck, theme_key: str | None = None) -> bytes:
    theme = get_theme(theme_key)
    prs = _new_presentation()
    expanded = _expand_animations(deck.slides)
    total = len(expanded)
    for i, slide in enumerate(expanded):
        _add_slide(prs, slide, theme, page_num=i + 1, total=total, lesson=deck.lesson_name)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- 内部实现 ---------------------------------------------------------------

def _new_presentation() -> Presentation:
    template = settings.template_pptx
    if template.exists():
        return Presentation(str(template))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _expand_animations(slides: list[Slide]) -> list[Slide]:
    out: list[Slide] = []
    for s in slides:
        if s.animation == AnimationHint.reveal_on_click and s.bullets:
            for i in range(1, len(s.bullets) + 1):
                out.append(s.model_copy(update={
                    "bullets": s.bullets[:i],
                    "animation": AnimationHint.none,
                }))
        elif s.animation == AnimationHint.step_by_step and s.solution_steps:
            for i in range(1, len(s.solution_steps) + 1):
                out.append(s.model_copy(update={
                    "solution_steps": s.solution_steps[:i],
                    "animation": AnimationHint.none,
                }))
        elif s.animation == AnimationHint.highlight_answer:
            out.append(s.model_copy(update={"answer": "", "animation": AnimationHint.none}))
            if s.answer:
                out.append(s.model_copy(update={"animation": AnimationHint.none}))
        else:
            out.append(s)
    return out


def _add_slide(prs: Presentation, slide: Slide, theme: Theme,
                page_num: int, total: int, lesson: str) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # 1. 整页背景（黑板主题）
    if theme.bg_color is not None:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.bg_color
        bg.line.fill.background()
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    has_diagram = bool(slide.diagram)

    example_sol_h = 0.0
    if slide.type == SlideType.title:
        _draw_title_slide(s, slide, theme)
    elif slide.type == SlideType.section:
        _draw_section_slide(s, slide, theme, page_num)
    elif slide.type == SlideType.example:
        example_sol_h = _draw_example_slide(s, slide, theme, compact=has_diagram) or 0.0
    elif slide.type == SlideType.practice:
        _draw_practice_slide(s, slide, theme, compact=has_diagram)
    elif slide.type == SlideType.interactive:
        _draw_interactive_slide(s, slide, theme, compact=has_diagram)
    elif slide.type == SlideType.summary:
        _draw_summary_slide(s, slide, theme)
    else:
        _draw_content_slide(s, slide, theme, compact=has_diagram)

    if has_diagram:
        if slide.type == SlideType.example:
            # diagram 跟在 solution 卡之后；y_top = sol_y + sol_h + gap
            diagram_top = 2.5 + max(example_sol_h, 1.5) + 0.1
            avail = 5.95 - diagram_top   # 6.05 是 answer 顶
            diagram_h = max(1.0, min(avail - 0.05, 1.85))
            _add_diagram(s, slide.diagram, theme,
                         top_inches=diagram_top,
                         height_inches=diagram_h, width_inches=9.0)
        else:
            _add_diagram(s, slide.diagram, theme)

    # 2. footer：除封面外都加
    if slide.type != SlideType.title:
        _add_footer(s, theme, page_num, total, lesson)

    if slide.notes or slide.duration_minutes:
        prefix = f"⏱ ~{slide.duration_minutes:.1f} 分钟\n\n" if slide.duration_minutes > 0 else ""
        s.notes_slide.notes_text_frame.text = prefix + (slide.notes or "")


# --- 各类型 slide 绘制 ------------------------------------------------------

def _draw_title_slide(s, slide: Slide, theme: Theme) -> None:
    """封面：左侧 5 寸大色块（带装饰）+ 右侧大标题 + chip 标签。"""
    # 左侧色块
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), SLIDE_H)
    block.fill.solid()
    block.fill.fore_color.rgb = theme.title_color
    block.line.fill.background()

    # 左侧大圆装饰
    deco1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(3.5), Inches(3.5))
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = theme.accent_color
    deco1.line.fill.background()
    deco1.shadow.inherit = False

    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(5), Inches(2.8), Inches(2.8))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = _lighten(theme.title_color, 0.3)
    deco2.line.fill.background()

    # 左下角"数学"标签（白色，叠在色块上）
    tag = _add_textbox(s, "数学", Inches(0.5), Inches(6.5), Inches(2), Inches(0.6),
                       size=20, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, font=theme.cn_font)

    # 右侧标题（大字）
    title_text = slide.title or "（未命名）"
    _add_textbox(s, title_text, Inches(5.0), Inches(2.4), Inches(8), Inches(1.8),
                 size=theme.base_size_cover_title, color=theme.title_color, bold=True,
                 font=theme.cn_font)

    # 副标题
    if slide.subtitle:
        _add_textbox(s, slide.subtitle, Inches(5.0), Inches(4.2), Inches(8), Inches(0.8),
                     size=24, color=theme.text_color, font=theme.cn_font)

    # chip：年级 / 学期 / 单元（从 subtitle 解析或留空）
    # 简化处理：用 deck.unit_name 等信息作为 chip 装饰


def _draw_section_slide(s, slide: Slide, theme: Theme, page_num: int) -> None:
    """章节页：大编号色块 + 居中标题。"""
    # 左侧大编号
    num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.5), Inches(2.5), Inches(2.5))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = theme.title_color
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.margin_top = Inches(0.3)
    tf.text = str(page_num)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(80)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = theme.cn_font

    # 标题（粗大）
    _add_textbox(s, slide.title, Inches(4.5), Inches(2.8), Inches(8.5), Inches(1.4),
                 size=48, color=theme.title_color, bold=True, font=theme.cn_font)

    # 装饰线
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.4),
                              Inches(2), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.accent_color
    line.line.fill.background()

    # bullets（如果有）作为本章子要点
    if slide.bullets:
        _add_bullets_block(s, slide.bullets, theme, x=Inches(4.5), y=Inches(4.8),
                            w=Inches(8.5), h=Inches(2),
                            size=22, color=theme.muted_color, prefix="› ")


def _draw_content_slide(s, slide: Slide, theme: Theme, compact: bool = False) -> None:
    """普通内容页：标题栏 + bullets（大字、纵向居中、左侧细色条）"""
    _add_title_bar(s, slide.title, theme)
    if not slide.bullets:
        return
    # 左侧小色条标识
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6),
                             Inches(0.1), Inches(5.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.accent_color
    bar.line.fill.background()

    # bullets：居中纵向铺开
    bottom = Inches(6.5) if compact else Inches(7.0)
    _add_bullets_block(s, slide.bullets, theme,
                       x=Inches(0.9), y=Inches(1.7),
                       w=Inches(12.0), h=bottom - Inches(1.7),
                       size=theme.base_size_bullet, color=theme.text_color, prefix="● ")


def _draw_example_slide(s, slide: Slide, theme: Theme, compact: bool = False) -> None:
    """例题：三个分明的色块区。

    精确布局（compact 即 has_diagram，y 单位 inch）：
      0.0~1.0     标题栏
      1.15~2.35   题目卡 (h=1.2)
      2.5~4.0     解答卡 (h=1.5, compact 字号 14)
      4.15~5.85   diagram (h=1.7, w=9, 居中)
      6.05~6.6    答案条 (h=0.55, 在 footer 7.05 上方)

    非 compact:
      1.15~2.35   题目
      2.5~6.0     解答（大）
      6.15~6.7    答案
    """
    _add_title_bar(s, slide.title or "例题", theme)

    # 题目卡片
    if slide.question:
        y = 1.15
        q_h = 1.2
        _add_card(s, theme,
                  left=Inches(0.5), top=Inches(y),
                  width=Inches(12.3), height=Inches(q_h),
                  bg=theme.soft_title)
        _add_textbox(s, "题目", Inches(0.7), Inches(y + 0.06), Inches(1.5), Inches(0.4),
                     size=16, color=theme.title_color, bold=True, font=theme.cn_font)
        _add_textbox(s, slide.question, Inches(0.7), Inches(y + 0.4),
                     Inches(11.9), Inches(q_h - 0.45),
                     size=22, color=theme.text_color, font=theme.cn_font)

    # 解答区（自适应高度 + 字号）
    sol_h = 0.0
    if slide.solution_steps:
        n = len(slide.solution_steps)
        total = sum(len(t) for t in slide.solution_steps)
        if compact:
            # 紧凑：步数多/字多时缩字号、增高度
            if n <= 3 and total <= 120:
                sol_h, step_size = 1.5, 14
            elif n <= 4 and total <= 200:
                sol_h, step_size = 1.85, 13
            else:
                sol_h, step_size = 2.2, 12
        else:
            sol_h = 3.5
            step_size = 18 if total > 200 else 20

        sol_y = 2.5
        _add_card(s, theme,
                  left=Inches(0.5), top=Inches(sol_y),
                  width=Inches(12.3), height=Inches(sol_h),
                  bg=None, border=theme.accent_color)
        _add_textbox(s, "解答", Inches(0.7), Inches(sol_y + 0.02),
                     Inches(1.5), Inches(0.4),
                     size=16, color=theme.accent_color, bold=True, font=theme.cn_font)
        _add_numbered_steps(s, slide.solution_steps, theme,
                             x=Inches(0.9), y=Inches(sol_y + 0.42),
                             w=Inches(11.6), h=Inches(sol_h - 0.5),
                             size=step_size)

    # 答案条
    if slide.answer:
        ans_y = 6.05 if compact else 6.15
        _add_card(s, theme,
                  left=Inches(0.5), top=Inches(ans_y),
                  width=Inches(12.3), height=Inches(0.55),
                  bg=theme.accent_color)
        _add_textbox(s, "✓  答：" + slide.answer,
                     Inches(0.8), Inches(ans_y + 0.05),
                     Inches(11.7), Inches(0.45),
                     size=20 if compact else 22,
                     color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     font=theme.cn_font)

    # 返回 sol_h 给 _add_slide 用以定位 diagram
    return sol_h


def _draw_practice_slide(s, slide: Slide, theme: Theme, compact: bool = False) -> None:
    """课堂练习：题目卡（自适应高度）+ 答题留白框 + 提示。"""
    _add_title_bar(s, slide.title or "课堂练习", theme)

    # 估算题目高度：按字符数 / 换行数粗算
    q = slide.question or ""
    line_breaks = q.count("\n\n") + 1
    rough_lines = max(line_breaks * 2, q.count("\n") + len(q) // 35 + 1)
    q_card_h = max(2.4, min(0.5 * rough_lines + 0.6, 4.5))

    if slide.question:
        _add_card(s, theme,
                  left=Inches(0.5), top=Inches(1.4),
                  width=Inches(12.3), height=Inches(q_card_h),
                  bg=theme.soft_accent, border=theme.accent_color)
        # 大问号图标
        _add_textbox(s, "?", Inches(0.7), Inches(1.5), Inches(0.8), Inches(0.8),
                     size=52, color=theme.accent_color, bold=True, font=theme.cn_font)
        _add_textbox(s, q, Inches(1.5), Inches(1.55),
                     Inches(11.1), Inches(q_card_h - 0.2),
                     size=22, color=theme.text_color, font=theme.cn_font)

    # 答题留白框：仅当题目卡占地不超过 4.6" 时才放
    ans_top = 1.4 + q_card_h + 0.2
    if not compact and ans_top + 1.0 < 6.7:
        ans_h = min(1.5, 6.5 - ans_top)
        _add_textbox(s, "你的解答：", Inches(0.5), Inches(ans_top), Inches(2.5), Inches(0.4),
                     size=18, color=theme.muted_color, font=theme.cn_font)
        ans_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.5), Inches(ans_top + 0.45),
                                      Inches(12.3), Inches(ans_h))
        ans_box.fill.background()
        ans_box.line.color.rgb = theme.muted_color
        ans_box.line.dash_style = 7
        ans_box.line.width = Pt(1.5)

    # hint 永远在 footer 上方
    if slide.hint and not compact:
        _add_textbox(s, "💡 提示  " + slide.hint, Inches(0.5), Inches(6.65),
                     Inches(12.3), Inches(0.35),
                     size=16, color=theme.muted_color, italic=True, font=theme.cn_font)


def _draw_interactive_slide(s, slide: Slide, theme: Theme, compact: bool = False) -> None:
    """互动页：大字居中问题 + 留白思考区 + hint。"""
    _add_title_bar(s, slide.title or "想一想", theme)

    # 中央大问题（高亮 chip）
    if slide.question:
        _add_card(s, theme,
                  left=Inches(0.8), top=Inches(2.0),
                  width=Inches(11.7), height=Inches(3.0),
                  bg=theme.soft_title, border=theme.title_color)
        # 思考图标
        _add_textbox(s, "💭", Inches(1.0), Inches(2.15), Inches(1), Inches(1),
                     size=44, color=theme.title_color, font=theme.cn_font)
        # 大字问题
        tb = _add_textbox(s, slide.question,
                          Inches(2.0), Inches(2.3),
                          Inches(10.5), Inches(2.4),
                          size=30, color=theme.title_color, bold=True,
                          font=theme.cn_font, align_center=True)

    # bullets：思考方向
    if slide.bullets and not compact:
        _add_bullets_block(s, slide.bullets, theme,
                            x=Inches(0.8), y=Inches(5.3),
                            w=Inches(11.7), h=Inches(1.4),
                            size=20, color=theme.text_color, prefix="› ")

    # hint
    if slide.hint and not compact:
        _add_textbox(s, "💡 提示  " + slide.hint, Inches(0.8), Inches(6.7),
                     Inches(11.7), Inches(0.5),
                     size=18, color=theme.muted_color, italic=True, font=theme.cn_font)


def _draw_summary_slide(s, slide: Slide, theme: Theme) -> None:
    """课堂小结：编号要点 + 强装饰。"""
    _add_title_bar(s, slide.title or "课堂小结", theme)

    # 大五星装饰
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(11.5), Inches(1.4), Inches(1.5), Inches(1.5))
    star.fill.solid()
    star.fill.fore_color.rgb = theme.accent_color
    star.line.fill.background()

    if not slide.bullets:
        return

    # 编号要点
    y = 1.7
    item_h = min(5.0 / max(len(slide.bullets), 1), 1.0)
    for i, b in enumerate(slide.bullets):
        # 编号圆
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.8), Inches(y),
                                     Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme.title_color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_top = Inches(0.05)
        tf.text = str(i + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = theme.cn_font

        # 要点文字
        _add_textbox(s, b, Inches(1.7), Inches(y + 0.08),
                     Inches(10.5), Inches(item_h),
                     size=24, color=theme.text_color, font=theme.cn_font)
        y += item_h


# --- 通用绘制 ---------------------------------------------------------------

def _add_title_bar(s, title: str, theme: Theme) -> None:
    """顶部标题栏 + 下方细色条，所有内容页通用。"""
    bar_shape = MSO_SHAPE.ROUNDED_RECTANGLE if theme.rounded_corners else MSO_SHAPE.RECTANGLE
    bar = s.shapes.add_shape(bar_shape, 0, 0, SLIDE_W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.title_color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = title or ""
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(theme.base_size_title_bar)
        run.font.bold = True
        run.font.color.rgb = theme.title_bar_text_color
        run.font.name = theme.cn_font

    # 下方装饰细线（accent 色）
    accent_line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      0, Inches(0.95), SLIDE_W, Inches(0.06))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = theme.accent_color
    accent_line.line.fill.background()


def _add_footer(s, theme: Theme, page_num: int, total: int, lesson: str) -> None:
    """底部 footer：左侧课题 / 右侧页码。"""
    # 横分隔线
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.muted_color
    line.line.fill.background()

    # 课题（左）
    _add_textbox(s, lesson, Inches(0.5), Inches(7.12), Inches(8), Inches(0.35),
                 size=12, color=theme.muted_color, font=theme.cn_font)
    # 页码（右）
    _add_textbox(s, f"{page_num} / {total}",
                 Inches(11.5), Inches(7.12), Inches(1.3), Inches(0.35),
                 size=12, color=theme.muted_color, font=theme.cn_font, align_right=True)


def _add_card(s, theme: Theme, left, top, width, height,
              bg: RGBColor | None = None, border: RGBColor | None = None) -> None:
    """背景色块卡片，做内容分区视觉。"""
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if theme.rounded_corners else MSO_SHAPE.RECTANGLE
    card = s.shapes.add_shape(shape, left, top, width, height)
    if bg is not None:
        card.fill.solid()
        card.fill.fore_color.rgb = bg
    else:
        card.fill.background()
    if border is not None:
        card.line.color.rgb = border
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()


def _add_textbox(s, text: str, left, top, width, height,
                  size: int = 20, bold: bool = False, italic: bool = False,
                  color: RGBColor | None = None, font: str = "Microsoft YaHei",
                  align_center: bool = False, align_right: bool = False):
    box = s.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = text
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    elif align_right:
        p.alignment = PP_ALIGN.RIGHT
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
        run.font.name = font
    return box


def _add_bullets_block(s, bullets: list[str], theme: Theme,
                        x, y, w, h, size: int = 28,
                        color: RGBColor | None = None, prefix: str = "• ") -> None:
    """大字号、宽行距、纵向均匀分布的 bullets。"""
    color = color or theme.text_color
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = prefix + item
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = theme.cn_font
        p.space_after = Pt(14)
        p.line_spacing = 1.25


def _add_numbered_steps(s, steps: list[str], theme: Theme,
                         x, y, w, h, size: int = 20) -> None:
    """例题解答步骤：编号 + 文字。字号小时收紧行距，防溢出。"""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    # 小字号用更紧凑的间距
    if size <= 13:
        line_spacing = 1.05
        space_after = Pt(2)
    elif size <= 16:
        line_spacing = 1.15
        space_after = Pt(5)
    else:
        line_spacing = 1.2
        space_after = Pt(8)
    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {i + 1}.  {step}"
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = theme.text_color
            run.font.name = theme.cn_font
        p.space_after = space_after
        p.line_spacing = line_spacing


def _add_diagram(s, diagram: dict, theme: Theme,
                  top_inches: float = 4.2, height_inches: float = 2.6,
                  width_inches: float = 11.0) -> None:
    """在 slide 居中嵌入数学示意图。可指定位置和大小。"""
    try:
        png = diagrams.render_diagram_png_bytes(diagram)
    except Exception as e:
        if s.has_notes_slide:
            s.notes_slide.notes_text_frame.text += f"\n[diagram error] {e}"
        return

    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        tmp.write(png)
        tmp_path = tmp.name
    try:
        width = Inches(width_inches)
        height = Inches(height_inches)
        left = Inches((13.333 - width_inches) / 2)
        top = Inches(top_inches)
        s.shapes.add_picture(tmp_path, left, top, width=width, height=height)
    finally:
        Path(tmp_path).unlink(missing_ok=True)
