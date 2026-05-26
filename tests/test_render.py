"""不调 API，只测 JSON → pptx 渲染。"""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation

from app.models.slide import Deck
from app.services.render import render, render_to_bytes


FIXTURE = Path(__file__).parent / "fixtures" / "sample_deck.json"


@pytest.fixture
def deck() -> Deck:
    return Deck.model_validate(json.loads(FIXTURE.read_text(encoding="utf-8")))


def test_deck_loads(deck: Deck) -> None:
    assert deck.grade == 4
    assert deck.term == 2
    assert deck.deck_type == "lesson_plan"
    assert len(deck.slides) >= 5


def test_render_to_file(deck: Deck, tmp_path: Path) -> None:
    out = render(deck, out_path=tmp_path / "out.pptx")
    assert out.exists()
    assert out.stat().st_size > 5000   # 不能是空文件

    # 读回去验证：动画展开后页数 > 原 slide 数
    prs = Presentation(str(out))
    assert len(prs.slides) >= len(deck.slides)


def test_render_to_bytes(deck: Deck) -> None:
    data = render_to_bytes(deck)
    assert data[:2] == b"PK"   # pptx 是 zip
    assert len(data) > 5000


def test_animation_expansion(deck: Deck, tmp_path: Path) -> None:
    """reveal_on_click 和 step_by_step 应让总页数 > 原 slide 数。"""
    raw_count = len(deck.slides)
    expanded = 0
    for s in deck.slides:
        if s.animation == "reveal_on_click":
            expanded += len(s.bullets) - 1
        elif s.animation == "step_by_step":
            expanded += len(s.solution_steps) - 1
        elif s.animation == "highlight_answer" and s.answer:
            expanded += 1
    expected = raw_count + expanded

    out = render(deck, out_path=tmp_path / "anim.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == expected


def test_notes_persisted(deck: Deck, tmp_path: Path) -> None:
    """有 notes 的 slide 应把讲稿写到 PPT 备注栏。"""
    out = render(deck, out_path=tmp_path / "notes.pptx")
    prs = Presentation(str(out))
    notes_found = [s.notes_slide.notes_text_frame.text for s in prs.slides
                   if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()]
    assert any("白菜" in n or "答案" in n or "对齐" in n for n in notes_found)
