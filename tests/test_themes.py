"""测三种主题都能渲染出有效的 .pptx。"""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation

from app.models.slide import Deck
from app.services.render import THEMES, get_theme, render


@pytest.fixture
def deck() -> Deck:
    raw = json.loads((Path(__file__).parent / "fixtures" / "sample_deck.json").read_text(encoding="utf-8"))
    return Deck.model_validate(raw)


def test_three_themes_registered():
    assert set(THEMES.keys()) == {"formal_blue", "kids_warm", "blackboard"}


def test_get_theme_fallback():
    assert get_theme(None).key == "formal_blue"
    assert get_theme("nonexistent").key == "formal_blue"
    assert get_theme("kids_warm").key == "kids_warm"


@pytest.mark.parametrize("theme_key", ["formal_blue", "kids_warm", "blackboard"])
def test_each_theme_renders(deck: Deck, theme_key: str, tmp_path: Path):
    out = render(deck, out_path=tmp_path / f"{theme_key}.pptx", theme_key=theme_key)
    assert out.exists()
    assert out.stat().st_size > 5000
    prs = Presentation(str(out))
    assert len(prs.slides) > 0


def test_blackboard_has_background_shape(deck: Deck, tmp_path: Path):
    """黑板风每页应有一个全屏背景矩形。"""
    out = render(deck, out_path=tmp_path / "bb.pptx", theme_key="blackboard")
    prs = Presentation(str(out))
    first = prs.slides[0]
    # 黑板风第一页至少有一个 shape 是 13.333x7.5 inches 的背景
    from pptx.util import Inches
    has_bg = any(
        abs(sh.width - Inches(13.333)) < Inches(0.1) and
        abs(sh.height - Inches(7.5)) < Inches(0.1)
        for sh in first.shapes
    )
    assert has_bg, "blackboard 主题首页应有全屏背景矩形"
